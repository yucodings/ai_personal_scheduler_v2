from io import BytesIO
from pptx import Presentation

def parse_pptx(content: bytes) -> str:
    presentation = Presentation(BytesIO(content)); slides = []
    for index, slide in enumerate(presentation.slides, 1):
        parts = [shape.text.strip() for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
        if parts: slides.append(f"[Slide {index}]\n" + "\n".join(parts))
    return "\n\n".join(slides)

